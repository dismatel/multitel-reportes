import azure.functions as func
import azure.durable_functions as df
import json
import logging
import os
import hashlib
import base64
import tempfile
import re
from datetime import datetime, timezone
from typing import Optional, List, Dict
from io import BytesIO

from ..shared.auth import require_auth, get_secret

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# SLIDE INDEX MAPS
# ---------------------------------------------------------------------------

PLANTA_EXTERNA_SLOTS: Dict[str, dict] = {
    "Punta Inicial":          {"slide": 2, "placeholder": 0},
    "Punta Final":            {"slide": 2, "placeholder": 1},
    "Reserva 1":              {"slide": 3, "placeholder": 0},
    "Reserva 2":              {"slide": 3, "placeholder": 1},
    "Reserva 3":              {"slide": 3, "placeholder": 2},
    "Reserva 4":              {"slide": 4, "placeholder": 0},
    "Reserva 5":              {"slide": 4, "placeholder": 1},
    "Cambio nodo 1":          {"slide": 5, "placeholder": 0},
    "Cambio nodo 2":          {"slide": 5, "placeholder": 1},
    "Servicios adicionales 1": {"slide": 5, "placeholder": 2},
    "Servicios adicionales 2": {"slide": 5, "placeholder": 3},
    "Metraje Punta Inicial":  {"slide": 2, "placeholder": 2},
    "Metraje Punta Final":    {"slide": 2, "placeholder": 3},
}

CPE_SLOTS: Dict[str, dict] = {
    "Rack sin CPE":     {"slide": 10, "placeholder": 0},
    "Rack con CPE":     {"slide": 10, "placeholder": 1},
    "Etiqueta LIU":     {"slide": 10, "placeholder": 2},
    "Etiqueta CPE":     {"slide": 11, "placeholder": 0},
    "Led Link":         {"slide": 11, "placeholder": 1},
    "ODF Nodo":         {"slide": 11, "placeholder": 2},
    "OLT/Switch":       {"slide": 11, "placeholder": 3},
    "Fusion Caja LIU":  {"slide": 12, "placeholder": 0},
    "Fusion Mufa":      {"slide": 13, "placeholder": 0},
    "SFP":              {"slide": 13, "placeholder": 1},
    "ADA":              {"slide": 14, "placeholder": 0},
    "ODI":              {"slide": 15, "placeholder": 0},
}

# ---------------------------------------------------------------------------
# RUN MERGING - CRITICAL: python-pptx splits {{Variable}} across multiple runs
# ---------------------------------------------------------------------------

def merge_runs_and_replace(paragraph, variables: dict) -> None:
    if not paragraph.runs:
        return
    full_text = "".join(run.text for run in paragraph.runs)
    replaced = False
    for key, value in variables.items():
        placeholder = f"{{{{{key}}}}}"
        if placeholder in full_text:
            full_text = full_text.replace(placeholder, str(value) if value is not None else "")
            replaced = True
    if replaced or any(f"{{{{{k}}}}}" in full_text for k in variables):
        paragraph.runs[0].text = full_text
        for run in paragraph.runs[1:]:
            run.text = ""


def replace_all_variables(prs, variables: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    merge_runs_and_replace(paragraph, variables)
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            merge_runs_and_replace(paragraph, variables)


# ---------------------------------------------------------------------------
# PHOTO INSERTION
# ---------------------------------------------------------------------------

def insert_photos(prs, fotos: list, tipo_reporte: str) -> None:
    from pptx.util import Inches
    from PIL import Image

    slot_map = PLANTA_EXTERNA_SLOTS if tipo_reporte == "Planta Externa" else CPE_SLOTS

    for foto in fotos:
        slot_nombre = foto.get("slot_nombre", "")
        imagen_b64 = foto.get("imagen_base64", "")
        if not imagen_b64 or slot_nombre not in slot_map:
            continue

        slot_info = slot_map[slot_nombre]
        slide_idx = slot_info["slide"]
        ph_idx = slot_info["placeholder"]

        if slide_idx >= len(prs.slides):
            logger.warning("slide_index_out_of_range slot=%s idx=%d", slot_nombre, slide_idx)
            continue

        slide = prs.slides[slide_idx]
        img_bytes = base64.b64decode(imagen_b64)

        target_ph = None
        picture_shapes = [sh for sh in slide.shapes
                          if hasattr(sh, "placeholder_format") and sh.placeholder_format is not None]
        if ph_idx < len(picture_shapes):
            target_ph = picture_shapes[ph_idx]

        if target_ph is not None:
            try:
                from pptx.util import Emu
                left = target_ph.left
                top = target_ph.top
                width = target_ph.width
                height = target_ph.height
                sp = target_ph._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(BytesIO(img_bytes), left, top, width, height)
            except Exception as exc:
                logger.warning("photo_insert_failed slot=%s err=%s", slot_nombre, exc)
        else:
            try:
                slide.shapes.add_picture(
                    BytesIO(img_bytes), Inches(1), Inches(1), Inches(4), Inches(3)
                )
            except Exception as exc:
                logger.warning("photo_append_failed slot=%s err=%s", slot_nombre, exc)


# ---------------------------------------------------------------------------
# BUILD VARIABLE MAP
# ---------------------------------------------------------------------------

def build_variable_map(payload: dict) -> dict:
    pc = payload.get("patchcord_vars", {})
    patchcord_map = {f"PC{i:02d}": pc.get(f"PC{i:02d}", "") for i in range(1, 29)}
    variables = {
        "Cliente":               payload.get("cliente", ""),
        "ID del Servicio":       payload.get("id_servicio", ""),
        "Encargado de grupo":    payload.get("encargado_grupo", ""),
        "Fecha":                 payload.get("fecha", ""),
        "Coordinadora":          payload.get("coordinadora", ""),
        "Encargados de grupos":  payload.get("encargados_grupos", ""),
        "Nodo":                  payload.get("nodo", ""),
        "Tipo de Servicio":      payload.get("tipo_servicio", ""),
        "Equipo Instalado":      payload.get("equipo_instalado", ""),
        "PotenciaCajaLiu":       payload.get("potencia_caja_liu", ""),
        "PerdidaCajaLiu":        payload.get("perdida_caja_liu", ""),
        "FusionCajaLiu":         payload.get("fusion_caja_liu", ""),
        "PerdidaMufaUltima":     payload.get("perdida_mufa_ultima", ""),
        "FusionMufaUlt":         payload.get("fusion_mufa_ult", ""),
        "InstSFP":               payload.get("inst_sfp", ""),
        "ODF":                   payload.get("odf", ""),
        "Rack sin CPE":          payload.get("rack_sin_cpe", ""),
        "Rack con CPE":          payload.get("rack_con_cpe", ""),
        "Etiqueta LIU":          payload.get("etiqueta_liu", ""),
        "Etiqueta CPE":          payload.get("etiqueta_cpe", ""),
        "Led Link":              payload.get("led_link", ""),
        "OLT/SWITCH":            payload.get("olt_switch", ""),
        "ODF Nodo":              payload.get("odf_nodo", ""),
        "ADA":                   payload.get("ada", ""),
        "ODI":                   payload.get("odi", ""),
        "Supervisor Lider":      payload.get("supervisor_lider", ""),
        "Firma Supervisor Lider": payload.get("firma_supervisor_lider", {}).get("imagen_base64", ""),
        "Gerente Operativo":     payload.get("gerente_operativo", ""),
        "Firma Gerente Operativo": payload.get("firma_gerente_operativo", {}).get("imagen_base64", ""),
        "Firma Coordinadora":    payload.get("firma_coordinadora", {}).get("imagen_base64", ""),
        "si/no":                 payload.get("si_no", ""),
        **patchcord_map,
    }
    return variables


# ---------------------------------------------------------------------------
# PDF CONVERSION - Python puro, sin LibreOffice
# ---------------------------------------------------------------------------

def convertir_pptx_a_pdf(pptx_bytes: bytes, reporte_id: str, tmp_dir: str) -> str:
    """
    Convierte PPTX a PDF usando reportlab + python-pptx.
    Sin dependencias de sistema operativo. Compatible con Azure Flex Consumption.
    """
    from pptx import Presentation
    from pptx.util import Pt
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
    from reportlab.lib import colors
    import PIL.Image as PILImage

    pdf_path = os.path.join(tmp_dir, f"reporte_{reporte_id}.pdf")
    prs = Presentation(BytesIO(pptx_bytes))

    doc = SimpleDocTemplate(
        pdf_path,
        pagesize=letter,
        rightMargin=inch * 0.75,
        leftMargin=inch * 0.75,
        topMargin=inch * 0.75,
        bottomMargin=inch * 0.75,
    )

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SlideTitle",
        parent=styles["Heading1"],
        fontSize=16,
        spaceAfter=6,
        textColor=colors.HexColor("#1F3864"),
    )
    body_style = ParagraphStyle(
        "SlideBody",
        parent=styles["Normal"],
        fontSize=11,
        spaceAfter=4,
        leading=15,
    )
    separator_style = ParagraphStyle(
        "Separator",
        parent=styles["Normal"],
        fontSize=8,
        textColor=colors.HexColor("#AAAAAA"),
        spaceAfter=10,
    )

    story = []
    slide_number = 0

    for slide in prs.slides:
        slide_number += 1
        has_content = False

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                has_content = True
                # Detectar si es titulo (primer parrafo del shape)
                try:
                    is_title = (
                        shape.placeholder_format is not None and
                        shape.placeholder_format.idx in (0, 1)
                    )
                except Exception:
                    is_title = False

                safe_text = (text
                    .replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;"))

                if is_title:
                    story.append(Paragraph(safe_text, title_style))
                else:
                    story.append(Paragraph(safe_text, body_style))

            # Insertar imágenes del shape si las hay
            try:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    img_stream = BytesIO(shape.image.blob)
                    pil_img = PILImage.open(img_stream)
                    w, h = pil_img.size
                    aspect = h / w
                    max_w = 5 * inch
                    img_w = min(max_w, w * 0.75)
                    img_h = img_w * aspect
                    img_stream.seek(0)
                    story.append(RLImage(img_stream, width=img_w, height=img_h))
                    story.append(Spacer(1, 6))
            except Exception as exc:
                logger.warning("pdf_image_skip slide=%d err=%s", slide_number, exc)

        if has_content and slide_number < len(prs.slides):
            story.append(Spacer(1, 12))
            story.append(Paragraph("─" * 60, separator_style))
            story.append(Spacer(1, 8))

    if not story:
        story.append(Paragraph("Reporte generado por Multitel", body_style))

    doc.build(story)
    logger.info("pdf_generated reporte_id=%s path=%s", reporte_id, pdf_path)
    return pdf_path


# ---------------------------------------------------------------------------
# DURABLE FUNCTION - HTTP STARTER
# ---------------------------------------------------------------------------

@require_auth(required_roles=["Tecnico"])
async def http_start(req: func.HttpRequest, starter: str, **kwargs) -> func.HttpResponse:
    client = df.DurableOrchestrationClient(starter)

    try:
        body = req.get_json()
    except ValueError:
        return func.HttpResponse(
            json.dumps({"error": "Payload JSON inválido"}),
            status_code=400,
            mimetype="application/json",
        )

    reporte_id = body.get("reporte_id")
    if not reporte_id:
        return func.HttpResponse(
            json.dumps({"error": "reporte_id requerido"}),
            status_code=400,
            mimetype="application/json",
        )

    _UUID_RE = re.compile(
        r'^[0-9a-f]{8}-[0-9a-f]{4}-4[0-9a-f]{3}-[89ab][0-9a-f]{3}-[0-9a-f]{12}$',
        re.IGNORECASE
    )
    if not _UUID_RE.match(reporte_id):
        logger.warning("fn_generar_pptx: reporte_id no es UUID válido: %s", reporte_id)
        return func.HttpResponse(
            json.dumps({"error": "reporte_id inválido"}),
            status_code=400,
            mimetype="application/json",
        )

    instance_id = await client.start_new("orchestrator", None, body)
    logger.info("durable_started reporte_id=%s instance=%s", reporte_id, instance_id)
    return client.create_check_status_response(req, instance_id)


# ---------------------------------------------------------------------------
# ORCHESTRATOR
# ---------------------------------------------------------------------------

def orchestrator(context: df.DurableOrchestrationContext):
    payload_data = context.get_input()
    reporte_id = payload_data.get("reporte_id")

    try:
        result = yield context.call_activity("generate_pptx_activity", payload_data)
        upload_input = {
            "reporte_id": reporte_id,
            "pptx_path": result["pptx_path"],
            "pdf_path": result["pdf_path"],
            "payload": payload_data.get("payload", {}),
        }
        upload_result = yield context.call_activity("upload_files_activity", upload_input)
        notify_input = {
            "reporte_id": reporte_id,
            "tecnico_email": payload_data.get("payload", {}).get("tecnico_email", ""),
            "pptx_url": upload_result.get("pptx_url", ""),
            "pdf_url": upload_result.get("pdf_url", ""),
        }
        yield context.call_activity("notify_activity", notify_input)
        return {"status": "completed", "reporte_id": reporte_id}

    except Exception as exc:
        logger.error("orchestrator_failed reporte_id=%s err=%s", reporte_id, exc)
        raise


# ---------------------------------------------------------------------------
# ACTIVITY: generate PPTX + convert to PDF
# ---------------------------------------------------------------------------

def generate_pptx_activity(payload_data: dict) -> dict:
    """
    1. Carga plantilla desde Azure Blob Storage
    2. Reemplaza variables
    3. Inserta fotos
    4. Guarda PPTX
    5. Convierte a PDF con reportlab (sin LibreOffice)
    6. Retorna rutas de ambos archivos
    """
    from pptx import Presentation
    from azure.storage.blob import BlobServiceClient

    reporte_id = payload_data.get("reporte_id")
    payload = payload_data.get("payload", {})
    tipo_reporte = payload.get("tipo_reporte", "Planta Externa")

    # Cargar plantilla desde Blob Storage
    blob_conn = get_secret("AZURE_STORAGE_CONNECTION_STRING")
    template_container = get_secret("TEMPLATE_CONTAINER_NAME")
    template_name = (
        "plantilla_planta_externa.pptx"
        if tipo_reporte == "Planta Externa"
        else "plantilla_cpe.pptx"
    )
    blob_client = BlobServiceClient.from_connection_string(blob_conn)
    container = blob_client.get_container_client(template_container)
    template_bytes = container.get_blob_client(template_name).download_blob().readall()

    prs = Presentation(BytesIO(template_bytes))

    # Reemplazar variables
    variables = build_variable_map(payload)
    replace_all_variables(prs, variables)

    # Insertar fotos
    fotos = payload.get("fotos", [])
    insert_photos(prs, fotos, tipo_reporte)

    # Guardar PPTX
    tmp_dir = tempfile.mkdtemp(prefix=f"reporte_{reporte_id}_")
    pptx_path = os.path.join(tmp_dir, f"reporte_{reporte_id}.pptx")
    prs.save(pptx_path)

    # SHA-256
    with open(pptx_path, "rb") as f:
        pptx_sha256 = hashlib.sha256(f.read()).hexdigest()
    logger.info("pptx_generated reporte_id=%s sha256=%s", reporte_id, pptx_sha256)

    try:
        _update_sha256_dataverse(reporte_id, pptx_sha256)
    except Exception as exc:
        logger.warning("sha256_update_failed reporte_id=%s err=%s", reporte_id, exc)

    # ---- CORRECCIÓN A-3: Convertir a PDF sin LibreOffice ----
    with open(pptx_path, "rb") as f:
        pptx_bytes = f.read()

    pdf_path = convertir_pptx_a_pdf(pptx_bytes, reporte_id, tmp_dir)

    return {
        "pptx_path": pptx_path,
        "pdf_path": pdf_path,
        "pptx_sha256": pptx_sha256,
    }


def _update_sha256_dataverse(reporte_id: str, sha256: str) -> None:
    import httpx
    from azure.identity import ManagedIdentityCredential, DefaultAzureCredential

    _mi_client_id = os.environ.get("MANAGED_IDENTITY_CLIENT_ID", "")
    cred = (
        ManagedIdentityCredential(client_id=_mi_client_id)
        if _mi_client_id
        else DefaultAzureCredential()
    )
    dataverse_url = get_secret("DATAVERSE_URL")
    token = cred.get_token(f"{dataverse_url.rstrip('/')}/.default").token
    httpx.patch(
        f"{dataverse_url}/api/data/v9.2/cr_multitelreportes"
        f"?$filter=cr_reporte_id eq '{reporte_id}'",
        headers={
            "Authorization": f"Bearer {token}",
            "Content-Type": "application/json",
            "OData-MaxVersion": "4.0",
            "OData-Version": "4.0",
        },
        json={"cr_pptx_sha256": sha256},
        timeout=15.0,
    )


# ---------------------------------------------------------------------------
# ACTIVITY: upload files
# ---------------------------------------------------------------------------

def upload_files_activity(upload_input: dict) -> dict:
    import httpx
    from azure.identity import ManagedIdentityCredential, DefaultAzureCredential

    reporte_id = upload_input["reporte_id"]
    pptx_path = upload_input["pptx_path"]
    pdf_path = upload_input["pdf_path"]

    _mi_client_id = os.environ.get("MANAGED_IDENTITY_CLIENT_ID", "")
    cred = (
        ManagedIdentityCredential(client_id=_mi_client_id)
        if _mi_client_id
        else DefaultAzureCredential()
    )
    graph_token = cred.get_token("https://graph.microsoft.com/.default").token
    drive_id = get_secret("ONEDRIVE_DRIVE_ID")
    base_path = f"/Multitel/Reportes/{reporte_id}"
    headers = {"Authorization": f"Bearer {graph_token}"}

    def upload_file(local_path: str, remote_name: str) -> str:
        with open(local_path, "rb") as f:
            content = f.read()
        url = (
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
            f"/root:{base_path}/{remote_name}:/content"
        )
        resp = httpx.put(
            url,
            headers={**headers, "Content-Type": "application/octet-stream"},
            content=content,
            timeout=120.0,
        )
        resp.raise_for_status()
        return resp.json().get("webUrl", "")

    pptx_url = upload_file(pptx_path, f"reporte_{reporte_id}.pptx")
    pdf_url = upload_file(pdf_path, f"reporte_{reporte_id}.pdf")
    logger.info("files_uploaded reporte_id=%s pptx=%s pdf=%s", reporte_id, pptx_url, pdf_url)

    try:
        import httpx as hx
        from azure.identity import DefaultAzureCredential as DAC
        cred2 = DAC()
        dv_url = get_secret("DATAVERSE_URL")
        dv_token = cred2.get_token(f"{dv_url.rstrip('/')}/.default").token
        hx.patch(
            f"{dv_url}/api/data/v9.2/cr_multitelreportes"
            f"?$filter=cr_reporte_id eq '{reporte_id}'",
            headers={
                "Authorization": f"Bearer {dv_token}",
                "Content-Type": "application/json",
                "OData-MaxVersion": "4.0",
                "OData-Version": "4.0",
            },
            json={"cr_pptx_url": pptx_url, "cr_pdf_url": pdf_url, "cr_estado": "Generado"},
            timeout=15.0,
        )
    except Exception as exc:
        logger.warning("dataverse_url_update_failed reporte_id=%s err=%s", reporte_id, exc)

    return {"pptx_url": pptx_url, "pdf_url": pdf_url}


# ---------------------------------------------------------------------------
# ACTIVITY: notify
# ---------------------------------------------------------------------------

def notify_activity(notify_input: dict) -> None:
    import httpx
    from azure.identity import ManagedIdentityCredential, DefaultAzureCredential

    reporte_id = notify_input.get("reporte_id", "")
    tecnico_email = notify_input.get("tecnico_email", "")
    pptx_url = notify_input.get("pptx_url", "")
    pdf_url = notify_input.get("pdf_url", "")

    _mi_client_id = os.environ.get("MANAGED_IDENTITY_CLIENT_ID", "")
    cred = (
        ManagedIdentityCredential(client_id=_mi_client_id)
        if _mi_client_id
        else DefaultAzureCredential()
    )
    graph_token = cred.get_token("https://graph.microsoft.com/.default").token
    pa_webhook = get_secret("POWER_AUTOMATE_APPROVAL_WEBHOOK")

    try:
        httpx.post(
            f"https://graph.microsoft.com/v1.0/users/{tecnico_email}/sendMail",
            headers={
                "Authorization": f"Bearer {graph_token}",
                "Content-Type": "application/json",
            },
            json={
                "message": {
                    "subject": f"Reporte {reporte_id} generado",
                    "body": {
                        "contentType": "HTML",
                        "content": (
                            f"<p>Tu reporte <b>{reporte_id}</b> fue generado exitosamente.</p>"
                            f'<p><a href="{pptx_url}">Descargar PPTX</a> | '
                            f'<a href="{pdf_url}">Descargar PDF</a></p>'
                        ),
                    },
                    "toRecipients": [{"emailAddress": {"address": tecnico_email}}],
                }
            },
            timeout=15.0,
        )
    except Exception as exc:
        logger.warning("email_notify_failed reporte_id=%s err=%s", reporte_id, exc)

    try:
        httpx.post(
            pa_webhook,
            json={
                "reporte_id": reporte_id,
                "tecnico_email": tecnico_email,
                "pptx_url": pptx_url,
                "pdf_url": pdf_url,
                "accion": "REPORTE_LISTO_PARA_APROBACION",
            },
            timeout=15.0,
        )
    except Exception as exc:
        logger.warning("teams_webhook_failed reporte_id=%s err=%s", reporte_id, exc)


# ---------------------------------------------------------------------------
# Durable Function app registration
# ---------------------------------------------------------------------------

main = df.Orchestrator.from_generator_function(orchestrator)
